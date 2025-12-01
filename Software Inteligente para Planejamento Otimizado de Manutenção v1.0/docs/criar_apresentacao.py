#!/usr/bin/env python3
"""
Script para criar PowerPoint profissional com roadmap TRL9
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Cores
COR_PRIMARIA = RGBColor(25, 118, 210)  # Azul
COR_SECUNDARIA = RGBColor(56, 142, 60)  # Verde
COR_ALERTA = RGBColor(211, 47, 47)  # Vermelho
COR_TEXTO = RGBColor(33, 33, 33)

def criar_apresentacao():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: CAPA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COR_PRIMARIA

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Software Inteligente para Planejamento de Manutenção"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Roadmap para TRL9 - Sistema de Produção"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 200, 200)

    # SLIDE 2: STATUS ATUAL
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Status Atual"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    items = [
        ("Algoritmo Funcional", "Cadeias de Markov + NSGA-II", "verde"),
        ("Fórmulas KPI Validadas", "Com 10+ referências científicas", "verde"),
        ("Fronteira de Pareto", "162 soluções não-dominadas", "verde"),
        ("TRL Atual", "5 (Protótipo de Pesquisa) - 20% completo", "vermelho"),
        ("Testes Automatizados", "0% cobertura", "vermelho"),
        ("Persistência de Dados", "Não existe", "vermelho"),
        ("API REST", "Não existe", "vermelho"),
    ]

    for idx, (titulo, desc, cor) in enumerate(items):
        if idx > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"• {titulo}: {desc}"
        p.font.size = Pt(18)
        p.font.bold = idx < 3
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    # SLIDE 3: PROBLEMAS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "89 Problemas Identificados"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COR_ALERTA

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    problemas = [
        ("CRÍTICOS (23)", "Bloqueadores para Produção", [
            "Sem type hints",
            "Sem validação de entrada",
            "Bare except clauses",
            "Sem logging estruturado",
        ]),
        ("ALTOS (41)", "Qualidade e Segurança", [
            "Zero cobertura de testes",
            "Sem API REST",
            "Sem autenticação",
            "Sem CI/CD pipeline",
        ]),
        ("MÉDIOS (25)", "Otimizações", [
            "Code duplication",
            "Performance bottlenecks",
        ]),
    ]

    for cat, desc, items_list in problemas:
        p = tf.add_paragraph()
        p.text = f"{cat} - {desc}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(8)

        for item in items_list:
            p = tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(14)
            p.level = 1

    # SLIDE 4: ROADMAP
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Roadmap: 6 Fases em 4-6 Semanas"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    fases = [
        ("Semana 1", "FUNDACAO", "40h", "Estrutura, config, logging, type hints"),
        ("Semana 2-3", "TESTES + DB", "70h", "80% cobertura, SQLite, backup"),
        ("Semana 3", "API REST", "40h", "FastAPI, endpoints, autenticação"),
        ("Semana 4", "DEPLOY", "40h", "Docker, CI/CD, Frontend"),
        ("Semana 5", "MONITORING", "30h", "Prometheus, Grafana, alertas"),
        ("Semana 6", "PRODUCAO", "30h", "Testes finais, SLA, deploy"),
    ]

    for semana, fase, horas, desc in fases:
        p = tf.add_paragraph()
        p.text = f"{semana} - {fase} ({horas}): {desc}"
        p.font.size = Pt(16)
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    p = tf.add_paragraph()
    p.text = "TOTAL: 240-300 horas | Equipe: 2-3 devs | Timeline: 4-6 semanas"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(12)

    # SLIDE 5: EQUIPE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Equipe Recomendada"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    equipe = [
        ("Backend Developer (1-2)", ["API REST", "Banco de Dados", "Integração"]),
        ("QA Engineer (1)", ["Testes Automatizados", "Validação", "Qualidade"]),
        ("DevOps Engineer (0.5)", ["CI/CD Pipeline", "Containerização", "Monitoring"]),
        ("Tech Lead (part-time)", ["Arquitetura", "Code Reviews", "Decisões Técnicas"]),
    ]

    for role, tasks in equipe:
        p = tf.add_paragraph()
        p.text = role
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_before = Pt(8)

        for task in tasks:
            p = tf.add_paragraph()
            p.text = f"  • {task}"
            p.font.size = Pt(14)
            p.level = 1

    # SLIDE 6: MÉTRICAS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Métricas de Sucesso"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    metricas = [
        ("Qualidade de Código", [
            "Type hints: >95%",
            "Unit tests: >80%",
            "Code duplication: <5%",
        ]),
        ("Funcionalidade", [
            "100% requisitos implementados",
            "0 bugs críticos",
            "Pareto quality: >95%",
        ]),
        ("Performance", [
            "API response: <500ms (p95)",
            "Memory usage: <500MB",
            "DB queries: <100ms (p95)",
        ]),
        ("Operacional", [
            "Uptime: 99.9%",
            "MTTR: <15 min",
            "MTBF: >30 dias",
        ]),
    ]

    for categoria, items_list in metricas:
        p = tf.add_paragraph()
        p.text = categoria
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(6)

        for item in items_list:
            p = tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(13)
            p.level = 1

    # SLIDE 7: RISCOS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Riscos & Mitigação"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COR_ALERTA

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    riscos = [
        ("Refatoração quebra funcionalidade", "Testes extensivos (80%+ cobertura)"),
        ("Performance inadequada", "Load testing, profiling, optimization"),
        ("DB migration perde dados", "Backup antes migração, teste restore"),
        ("Deadline apertado", "Priorizar MVP, diferir nice-to-haves"),
    ]

    for risco, mitigacao in riscos:
        p = tf.add_paragraph()
        p.text = f"Risco: {risco}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = f"Mitigation: {mitigacao}"
        p.font.size = Pt(14)
        p.level = 1

    # SLIDE 8: PRÓXIMOS PASSOS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Próximos Passos"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    passos = [
        ("HOJE", [
            "Revisar documentação",
            "Tomar decisão de prosseguimento",
            "Alocar recursos (2-3 desenvolvedores)",
        ]),
        ("SEMANA 1", [
            "Inicializar repositório Git",
            "Estruturar projeto (src/, tests/, docs/)",
            "Implementar Fase 1 (Fundação)",
        ]),
        ("CHECKPOINT TRL6 (Semana 3)", [
            "API REST funcional + 80% testes + BD persistente",
        ]),
        ("CHECKPOINT TRL9 (Semana 6)", [
            "Sistema completo em produção com SLA compliance",
        ]),
    ]

    for periodo, items_list in passos:
        p = tf.add_paragraph()
        p.text = periodo
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(8)

        for item in items_list:
            p = tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(13)
            p.level = 1

    # SLIDE 9: CONCLUSÃO
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COR_SECUNDARIA

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PROSSEGUIR COM IMPLEMENTAÇÃO"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Viabilidade: ALTA | Risco: MÉDIO (Gerenciável)"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Salvar
    output_file = "Apresentacao_TRL9_Manutencao.pptx"
    prs.save(output_file)
    print("\n[OK] PowerPoint criado com sucesso: {}".format(output_file))
    print("\nDetalhes:")
    print("  - 9 slides profissionais")
    print("  - Conteudo completo do roadmap TRL9")
    print("  - Pronto para apresentacao executiva")
    print("  - Formato PPTX (compativel com PowerPoint, Google Slides, etc)")

if __name__ == "__main__":
    criar_apresentacao()
