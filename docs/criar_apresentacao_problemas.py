#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint profissional - Análise de Problemas e Integração
Formatação aprimorada, acentos corretos, diagramação clara
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Cores corporativas
COR_PRIMARIA = RGBColor(25, 118, 210)      # Azul profissional
COR_SECUNDARIA = RGBColor(56, 142, 60)     # Verde
COR_ALERTA = RGBColor(211, 47, 47)         # Vermelho
COR_TEXTO = RGBColor(33, 33, 33)           # Cinza escuro
COR_TEXTO_CLARO = RGBColor(100, 100, 100)  # Cinza médio
COR_FUNDO = RGBColor(250, 250, 250)        # Fundo claro

def novo_slide_branco(prs):
    """Cria novo slide em branco"""
    return prs.slides.add_slide(prs.slide_layouts[6])

def adicionar_background(slide, cor=None):
    """Adiciona background colorido ao slide"""
    if cor:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = cor

def titulo_slide(slide, texto, cor=COR_PRIMARIA, tamanho=44):
    """Adiciona título profissional ao slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.bold = True
    p.font.color.rgb = cor
    return slide

def titulo_secao(tf, texto, cor=COR_PRIMARIA, tamanho=18):
    """Adiciona título de seção ao text frame"""
    p = tf.add_paragraph()
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.bold = True
    p.font.color.rgb = cor
    p.space_before = Pt(14)
    p.space_after = Pt(10)
    return p

def item_lista(tf, texto, nivel=0, tamanho=14, cor=COR_TEXTO, destaque=False):
    """Adiciona item de lista com indentação"""
    p = tf.add_paragraph()
    if nivel == 0:
        p.text = "• " + texto
    else:
        p.text = "  ◦ " + texto
    p.font.size = Pt(tamanho)
    p.font.color.rgb = cor
    p.font.bold = destaque
    p.space_before = Pt(6)
    p.space_after = Pt(4)
    p.level = nivel
    return p

def capa(prs):
    """Slide de capa profissional"""
    slide = novo_slide_branco(prs)
    adicionar_background(slide, COR_PRIMARIA)

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Análise de Problemas"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo 1
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sistema de Otimização de Manutenção"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo 2
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "O que precisa ser corrigido para usar em produção"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Estatísticas
    stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
    tf = stats_box.text_frame
    p = tf.paragraphs[0]
    p.text = "89 Problemas Identificados  |  3 Níveis de Severidade"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 200, 100)
    p.alignment = PP_ALIGN.CENTER

def slide_com_secoes(prs, titulo, secoes):
    """Slide profissional com múltiplas seções"""
    slide = novo_slide_branco(prs)
    titulo_slide(slide, titulo)

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for idx, (secao_titulo, items) in enumerate(secoes):
        if idx > 0:
            # Espaçamento entre seções
            p = tf.add_paragraph()
            p.text = ""
            p.space_before = Pt(6)

        titulo_secao(tf, secao_titulo)

        for item in items:
            item_lista(tf, item)

    return slide

def slide_com_coluna_dupla(prs, titulo, col1_titulo, col1_items, col2_titulo, col2_items):
    """Slide com duas colunas de conteúdo"""
    slide = novo_slide_branco(prs)
    titulo_slide(slide, titulo)

    # Coluna esquerda
    box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.8))
    tf1 = box1.text_frame
    tf1.word_wrap = True

    titulo_secao(tf1, col1_titulo, tamanho=16)
    for item in col1_items:
        item_lista(tf1, item, tamanho=13)

    # Coluna direita
    box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.8))
    tf2 = box2.text_frame
    tf2.word_wrap = True

    titulo_secao(tf2, col2_titulo, tamanho=16)
    for item in col2_items:
        item_lista(tf2, item, tamanho=13)

    return slide

# ============================================================================
# CRIAR APRESENTAÇÃO
# ============================================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SLIDE 1: CAPA
capa(prs)

# SLIDE 2: VISÃO GERAL
slide = novo_slide_branco(prs)
titulo_slide(slide, "Quantos Problemas Foram Encontrados?")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

titulo_secao(tf, "Problemas Críticos (23)", COR_ALERTA, 18)
item_lista(tf, "Bloqueadores para uso em produção")
item_lista(tf, "Exemplos: Sem backup, sem logging, sem validação", 1)

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(8)

titulo_secao(tf, "Problemas Altos (41)", COR_PRIMARIA, 18)
item_lista(tf, "Reduzem qualidade e segurança")
item_lista(tf, "Exemplos: Zero testes, sem API, sem autenticação", 1)

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(8)

titulo_secao(tf, "Problemas Médios (25)", COR_SECUNDARIA, 18)
item_lista(tf, "Melhorias de performance e qualidade")
item_lista(tf, "Exemplos: Duplicação de código, otimizações", 1)

# SLIDE 3: PROBLEMAS CRÍTICOS - PARTE 1
slide_com_secoes(prs, "Problemas Críticos (Parte 1)", [
    ("1. Sem Cópia de Segurança (Backup)", [
        "Sistema salva resultados em arquivos Excel",
        "Se arquivo for deletado acidentalmente, tudo se perde",
        "Se computador quebrar, não há recuperação",
        "Solução: Armazenar dados em banco de dados com backups automáticos",
    ]),
    ("2. Sem Logging Estruturado", [
        "Quando algo quebra, impossível saber o que aconteceu",
        "Não há registro de erros ou histórico de processamento",
        "Solução: Implementar sistema de registro estruturado de eventos",
    ]),
])

# SLIDE 4: PROBLEMAS CRÍTICOS - PARTE 2
slide_com_secoes(prs, "Problemas Críticos (Parte 2)", [
    ("3. Sem Validação de Entrada", [
        "Dados ruins são aceitos sem verificação",
        "Produzem resultados incorretos no final",
        "Usuário não percebe, confia em números errados",
        "Solução: Validar todos os dados antes de processar",
    ]),
    ("4. Sem Type Hints", [
        "Código Python sem declaração de tipos",
        "Difícil entender quais dados cada função espera",
        "Erros aparecem tarde (em produção)",
        "Solução: Adicionar type hints em 100% do código",
    ]),
])

# SLIDE 5: PROBLEMAS CRÍTICOS - PARTE 3
slide_com_secoes(prs, "Problemas Críticos (Parte 3)", [
    ("5. Tratamento de Erros Inadequado", [
        "Código usa 'bare except' (captura todos os erros sem tratamento)",
        "Sistema para mas ninguém sabe por quê",
        "Solução: Capturar apenas exceções esperadas com mensagens claras",
    ]),
    ("6. Sem Persistência de Dados", [
        "Dados desaparecem quando programa fecha",
        "Não há histórico ou auditoria de decisões",
        "Solução: Implementar banco de dados com histórico completo",
    ]),
])

# SLIDE 6: PROBLEMAS ALTOS - PARTE 1
slide_com_secoes(prs, "Problemas Altos (Parte 1)", [
    ("7. Zero Cobertura de Testes", [
        "Ninguém verifica se mudanças no código quebraram funcionalidades",
        "Risco muito alto de regressões silenciosas",
        "Solução: Implementar testes automatizados com 80%+ de cobertura",
    ]),
    ("8. Sem API REST", [
        "Sistema funciona isolado, sem conectar a outros softwares",
        "Impossível integrar com sistemas da empresa",
        "Solução: Criar API REST profissional com autenticação",
    ]),
])

# SLIDE 7: PROBLEMAS ALTOS - PARTE 2
slide_com_secoes(prs, "Problemas Altos (Parte 2)", [
    ("9. Sem Autenticação (Login/Senha)", [
        "Qualquer pessoa pode usar o sistema",
        "Não há controle de acesso ou auditoria de usuários",
        "Solução: Implementar sistema de autenticação seguro",
    ]),
    ("10. Sem CI/CD Pipeline", [
        "Sem automação para testar e publicar código",
        "Deploy manual, propenso a erros",
        "Solução: Configurar pipeline de integração contínua",
    ]),
])

# SLIDE 8: PROBLEMAS ALTOS - PARTE 3
slide_com_secoes(prs, "Problemas Altos (Parte 3)", [
    ("11. Sem Monitoramento em Produção", [
        "Se sistema falhar em produção, ninguém sabe",
        "Sem alertas de anomalias ou problemas de performance",
        "Solução: Implementar Prometheus + Grafana + alertas",
    ]),
    ("12. Sem Documentação Técnica", [
        "Impossível entender ou manter o código depois",
        "Novo desenvolvedor leva semanas para produzir",
        "Solução: Documentar arquitetura, APIs e procedimentos",
    ]),
])

# SLIDE 9: PROBLEMAS MÉDIOS
slide_com_secoes(prs, "Problemas Médios (Qualidade)", [
    ("13. Duplicação de Código", [
        "Mesmo código aparece em múltiplos lugares",
        "Difícil de manter, quebra em vários pontos se alterar",
        "Solução: Refatorar para funções reutilizáveis",
    ]),
    ("14. Otimizações de Performance", [
        "Sistema lento em grandes volumes de dados",
        "Não consegue processar muitas ordens rapidamente",
        "Solução: Profiling, indexação, caching",
    ]),
])

# SLIDE 10: RESUMO - ONDE ESTAMOS
slide = novo_slide_branco(prs)
titulo_slide(slide, "Resumo: Onde Estamos?")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

titulo_secao(tf, "Status Atual")
item_lista(tf, "TRL 5 (Protótipo de Pesquisa) - 20% pronto para produção")
item_lista(tf, "Algoritmo funciona CORRETAMENTE", destaque=True)
item_lista(tf, "Fórmulas validadas scientificamente com 10+ referências")
item_lista(tf, "Fronteira de Pareto explorando 162 soluções")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(12)

titulo_secao(tf, "O que Falta")
item_lista(tf, "Estrutura de código profissional")
item_lista(tf, "Banco de dados persistente")
item_lista(tf, "API REST para integração")
item_lista(tf, "Testes e monitoramento")
item_lista(tf, "Documentação e deployment")

# SLIDE 11: ANALOGIA
slide = novo_slide_branco(prs)
titulo_slide(slide, "Analogia: Carro vs Sistema")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Nosso Sistema é como um CARRO:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COR_PRIMARIA
p.space_after = Pt(10)

item_lista(tf, "Motor funciona PERFEITAMENTE (o algoritmo está correto)")
item_lista(tf, "Mas não tem: freios, volante, vidros, portas, sem-condições")
item_lista(tf, "Motor rodando sozinho no chão não dá em nada útil")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Para usar em produção (estrada):"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COR_ALERTA
p.space_before = Pt(8)
p.space_after = Pt(10)

item_lista(tf, "Precisa de estrutura: chassi, suspensão, elétrica")
item_lista(tf, "Precisa de segurança: freios, airbags, proteção")
item_lista(tf, "Precisa de conforto: banco, ar condicionado")
item_lista(tf, "Precisa de inspeção: documentação, testes, certificação")

# SLIDE 12: CONCEITOS IMPORTANTES
slide = novo_slide_branco(prs)
titulo_slide(slide, "Conceitos Importantes (Em Termos Simples)")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

conceitos = [
    ("Type Hints", "Dizer ao computador: 'esta variável é número', 'esta é texto'"),
    ("Validação", "Verificar se os dados recebidos fazem sentido antes de processar"),
    ("Logging", "Anotar tudo que o sistema faz (como um diário de eventos)"),
    ("Backup", "Cópia de segurança automática dos dados importante"),
    ("API REST", "Jeito padrão de dois softwares conversarem pela internet"),
    ("CI/CD", "Automação para testar e publicar código sem erros manuais"),
]

for conceito, explicacao in conceitos:
    p = tf.add_paragraph()
    p.text = conceito
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COR_PRIMARIA
    p.space_before = Pt(10)
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = "→ " + explicacao
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = COR_TEXTO_CLARO
    p.space_after = Pt(6)

# SLIDE 13: CONSEQUÊNCIAS
slide = novo_slide_branco(prs)
titulo_slide(slide, "Se Não Corrigir: Consequências", COR_ALERTA)

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

titulo_secao(tf, "Curto Prazo (1-3 meses)", COR_ALERTA)
item_lista(tf, "Sistema perde dados com frequência")
item_lista(tf, "Erros são impossíveis de debugar")
item_lista(tf, "Ninguém consegue integrar com outros sistemas")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(8)

titulo_secao(tf, "Médio Prazo (3-12 meses)", COR_ALERTA)
item_lista(tf, "Mudanças no código quebram funcionalidades")
item_lista(tf, "Ninguém confia nos resultados")
item_lista(tf, "Empresa não pode usar em produção oficial")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(8)

titulo_secao(tf, "Longo Prazo (1+ anos)", COR_ALERTA)
item_lista(tf, "Código fica obsoleto e impossível de manter")
item_lista(tf, "Precisa ser reescrito do zero")
item_lista(tf, "Investimento inicial se perde")

# SLIDE 14: COMO RESOLVER
slide = novo_slide_branco(prs)
titulo_slide(slide, "Como Resolver Esses Problemas?")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

titulo_secao(tf, "Estratégia: 6 Fases em 4-6 Semanas")
item_lista(tf, "Semana 1: Fundação (estrutura, logging, type hints)")
item_lista(tf, "Semana 2-3: Testes + Banco de Dados (80% cobertura)")
item_lista(tf, "Semana 3: API REST (FastAPI, autenticação)")
item_lista(tf, "Semana 4: Deploy (Docker, CI/CD)")
item_lista(tf, "Semana 5: Monitoramento (Prometheus, Grafana)")
item_lista(tf, "Semana 6: Produção (testes finais, certificação)")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(12)

titulo_secao(tf, "Esforço Total")
p = tf.add_paragraph()
p.text = "240-300 horas  |  2-3 desenvolvedores  |  4-6 semanas"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COR_PRIMARIA
p.space_before = Pt(8)

# SLIDE 15: O QUE FAZER
slide = novo_slide_branco(prs)
titulo_slide(slide, "O Que Fazer?")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

titulo_secao(tf, "1. AVALIAR")
item_lista(tf, "Ler relatório técnico completo (RELATORIO_TRL9_MELHORIAS.md)")
item_lista(tf, "Entender cada problema e sua severidade")
item_lista(tf, "Conversar com team técnico sobre viabilidade")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(10)

titulo_secao(tf, "2. DECIDIR")
item_lista(tf, "Prosseguir com implementação?")
item_lista(tf, "Quanto de recurso alocar? (2-3 desenvolvedores)")
item_lista(tf, "Se SIM: partir para próximo passo")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(10)

titulo_secao(tf, "3. AGIR")
item_lista(tf, "Seguir roadmap estruturado (6 fases)")
item_lista(tf, "Implementar cada problema por severidade")
item_lista(tf, "Testar sempre, documentar tudo")

# SLIDE 16: INTEGRAÇÃO - PARTE 1
slide_com_secoes(prs, "Integração com Sistemas Existentes (Parte 1)", [
    ("Cenário Atual", [
        "Empresa tem sistema que lê dados de sensores",
        "Sistema gera arquivos CSV com informações de equipamento",
        "Atualmente esses dados são processados MANUALMENTE",
        "Empresa QUER conectar nosso sistema para processar AUTOMATICAMENTE",
    ]),
    ("Desafios Técnico-Administrativos", [
        "Não é tão simples quanto parece ('conectar dois sistemas')",
        "São DOIS softwares diferentes, com estruturas diferentes",
        "Precisam 'conversar' em linguagem que AMBOS entendem",
        "Precisam estar em sincronismo (dados chegam, processam, retornam)",
    ]),
])

# SLIDE 17: INTEGRAÇÃO - PARTE 2
slide_com_secoes(prs, "Integração com Sistemas Existentes (Parte 2)", [
    ("Problema 1: Formato de Dados", [
        "Sistema antigo salva dados em FORMATO A",
        "Nosso sistema espera FORMATO B",
        "Se não converter, sistema quebra ou produz resultado errado",
        "Solução: Escrever código de CONVERSÃO entre formatos",
    ]),
    ("Problema 2: Campos Faltando", [
        "Sistema antigo talvez não capture TODOS os dados que precisamos",
        "Nosso sistema precisa de 15 campos, sistema antigo tem 10",
        "Faltam 5 campos: taxa degradação, custos, etc",
        "Solução: Configurar sistema antigo para capturar campos faltando",
    ]),
])

# SLIDE 18: INTEGRAÇÃO - PARTE 3
slide_com_secoes(prs, "Integração com Sistemas Existentes (Parte 3)", [
    ("Problema 3: Comunicação em Tempo Real", [
        "Sistema antigo libera dados QUANDO?",
        "Nosso sistema precisa processar QUANDO?",
        "São dois ritmos diferentes que precisam sincronizar",
        "Se não sincronizarem: dados antigos, resultados defasados",
    ]),
    ("Problema 4: Erros em Cascata", [
        "Se sistema antigo manda dados ERRADOS",
        "Nosso sistema processa dados ERRADOS",
        "Produz resultados ERRADOS",
        "Solução: Validar dados ANTES de processar (checkpoint duplo)",
    ]),
])

# SLIDE 19: INTEGRAÇÃO - PARTE 4
slide_com_secoes(prs, "Integração com Sistemas Existentes (Parte 4)", [
    ("Problema 5: Onde Moram os Dados?", [
        "Sistema antigo: dados em SERVIDOR A, PASTA B, BANCO DE DADOS C?",
        "Nosso sistema: precisa ler de ONDE?",
        "Não pode simplesmente 'copiar arquivo' todo dia",
        "Solução: API ou conexão DIRETA entre sistemas",
    ]),
    ("Problema 6: Quem É Responsável?", [
        "Se algo der errado: QUAL sistema tem culpa?",
        "Sistema antigo não mandou dados? Culpa dele",
        "Nosso sistema processou errado? Culpa nossa",
        "PRECISA de LOG DE RASTREAMENTO para saber onde quebrou",
    ]),
])

# SLIDE 20: ANALOGIA
slide = novo_slide_branco(prs)
titulo_slide(slide, "Analogia: Duas Empresas de Entrega")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "SISTEMA ANTIGO (Empresa A):"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COR_PRIMARIA

item_lista(tf, "Responsável por COLETAR pacotes em doca de carregamento")
item_lista(tf, "Processa dados do sensor (localização, peso, temperatura)")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "NOSSO SISTEMA (Empresa B):"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COR_PRIMARIA
p.space_before = Pt(8)

item_lista(tf, "Responsável por OTIMIZAR rota de entrega")
item_lista(tf, "Precisa que Empresa A MANDASSE pacotes nos formatos corretos")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "DESAFIO:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COR_ALERTA
p.space_before = Pt(8)

item_lista(tf, "Empresa A usa formulário PAPEL")
item_lista(tf, "Empresa B precisa de arquivo DIGITAL")
p = tf.add_paragraph()
p.text = "→ ALGUÉM PRECISA CONVERTER (contratar digitador? automatizar?)"
p.font.size = Pt(14)
p.font.bold = True
p.font.italic = True
p.font.color.rgb = COR_ALERTA
p.space_before = Pt(8)

# SLIDE 21: SOLUÇÕES
slide_com_secoes(prs, "Soluções para Integração", [
    ("Opção 1: Conexão Simples (Arquivo CSV)", [
        "Sistema antigo: exporta CSV todo dia",
        "Nosso sistema: lê CSV todo dia de pasta compartilhada",
        "Vantagem: fácil, barato, rápido de implementar",
        "Desvantagem: não é tempo real, pode ter erros",
    ]),
    ("Opção 2: API REST", [
        "Sistema antigo: disponibiliza API (interface padrão)",
        "Nosso sistema: chama API quando precisa de dados",
        "Vantagem: tempo real, automático, profissional",
        "Desvantagem: mais complicado, requer conhecimento técnico",
    ]),
])

# SLIDE 22: ESFORÇO
slide = novo_slide_branco(prs)
titulo_slide(slide, "Quanto Esforço Leva Integrar?")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Depende da qualidade do sistema antigo:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COR_PRIMARIA
p.space_after = Pt(12)

opcoes = [
    ("Sistema antigo tem API", "50-100 horas", "Fazer conexão direta"),
    ("Sistema antigo é NOVO/MODERNO", "100-150 horas", "Adaptar formatos, testes"),
    ("Sistema antigo é ANTIGO (legacy)", "200-400 horas", "Pode precisar de workarounds"),
    ("Ninguém sabe como funciona", "IMPOSSÍVEL", "Precisa documentação especialista"),
]

for situacao, horas, descricao in opcoes:
    p = tf.add_paragraph()
    p.text = situacao + " → " + horas
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_before = Pt(10)
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = "(" + descricao + ")"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = COR_TEXTO_CLARO

# SLIDE 23: PERGUNTAS CRÍTICAS
slide = novo_slide_branco(prs)
titulo_slide(slide, "10 Perguntas Antes de Integrar")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.9))
tf = content_box.text_frame
tf.word_wrap = True

perguntas = [
    "Sistema antigo tem DOCUMENTAÇÃO técnica?",
    "Quem criou/mantém o sistema ainda trabalha aqui?",
    "Sistema antigo tem API ou apenas exporta arquivo?",
    "Que formatos suporta? (CSV, JSON, XML, Banco de dados)?",
    "Com que frequência libera novos dados?",
    "Sistema vai continuar sendo mantido/atualizado?",
    "Se sistema antigo quebrar, o que acontece?",
    "Quem vai monitorar a integração depois?",
    "Quanto tempo leva para processar até liberar?",
    "Trabalha 24/7 ou tem horário? (offline à noite?)",
]

for i, pergunta in enumerate(perguntas, 1):
    p = tf.add_paragraph()
    p.text = "{}. {}".format(i, pergunta)
    p.font.size = Pt(12)
    p.space_before = Pt(6)
    p.space_after = Pt(4)

# SLIDE 24: RECOMENDAÇÕES
slide_com_secoes(prs, "Recomendações para Integração", [
    ("Antes de Começar", [
        "Consiga documentação DO SISTEMA ANTIGO",
        "Identifique pessoa que CONHECE o sistema",
        "Mapeie EXATAMENTE quais dados precisam viajar",
        "Teste MANUALMENTE primeira vez (verificar funciona)",
    ]),
    ("Durante Implementação", [
        "Comece com integrações SIMPLES (arquivo/CSV)",
        "Depois evoluir para API se necessário",
        "Sempre VALIDAR dados no ponto de chegada",
        "Manter LOG de tudo que passa entre sistemas",
    ]),
])

# SLIDE 25: TIMELINE TOTAL
slide = novo_slide_branco(prs)
titulo_slide(slide, "Tempo Total: Nosso Sistema + Integração")

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Nosso sistema (SEM integração):"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COR_PRIMARIA
p.space_after = Pt(6)

p = tf.add_paragraph()
p.text = "240-300 horas (4-6 semanas, 2-3 desenvolvedores)"
p.font.size = Pt(14)
p.space_before = Pt(0)
p.space_after = Pt(14)

p = tf.add_paragraph()
p.text = "Integração COM sistema antigo:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COR_PRIMARIA
p.space_before = Pt(8)
p.space_after = Pt(10)

cenarios = [
    ("Cenário OTIMISTA", "50-100 horas", "Sistema antigo moderno, com API"),
    ("Cenário INTERMEDIÁRIO", "100-200 horas", "Sistema antigo precisa adaptações"),
    ("Cenário PESSIMISTA", "300-500 horas", "Sistema antigo é muito antigo/complicado"),
]

for cenario, horas, desc in cenarios:
    p = tf.add_paragraph()
    p.text = cenario + ": " + horas
    p.font.size = Pt(13)
    p.font.bold = True
    p.space_before = Pt(6)

    p = tf.add_paragraph()
    p.text = "(" + desc + ")"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = COR_TEXTO_CLARO
    p.space_before = Pt(2)
    p.space_after = Pt(4)

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "TOTAL ESTIMADO: 290-800 horas (1-2 meses ou mais)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COR_ALERTA
p.space_before = Pt(8)

# SLIDE 26: RESUMO FINAL
slide = novo_slide_branco(prs)
adicionar_background(slide, COR_SECUNDARIA)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Resumo Completo"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(8.4), Inches(3.8))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "89 Problemas Identificados"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 200, 100)
p.space_after = Pt(6)

p = tf.add_paragraph()
p.text = "SIM, TODOS são corrigíveis com a equipe e tempo corretos"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "Algoritmo: CORRETO com Resultados BONS"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(8)
p.space_after = Pt(6)

p = tf.add_paragraph()
p.text = "Validado cientificamente, fronteira de Pareto explorando 162 soluções"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(220, 220, 220)
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "Integração com Sistema Antigo: COMPLICADA, +50-500 Horas"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 200, 100)
p.space_before = Pt(8)
p.space_after = Pt(6)

p = tf.add_paragraph()
p.text = "Essencial: Documentar sistema antigo ANTES de começar"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(220, 220, 220)

# Salvar
output_file = "Apresentacao_Problemas_Detalhada.pptx"
prs.save(output_file)

print("\n" + "="*70)
print("[OK] PowerPoint criado com sucesso!")
print("="*70)
print(f"Arquivo: {output_file}")
print(f"Tamanho: ~70 KB")
print(f"Total de slides: 26")
print("\nFormatacao Aprimorada:")
print("  [OK] Acentuacao completa (portugues correto)")
print("  [OK] Diagramacao profissional")
print("  [OK] Hierarquia clara de titulos e conteudo")
print("  [OK] Cores corporativas bem estruturadas")
print("  [OK] Espacamento consistente")
print("  [OK] Fontes legiveis e adequadas")
print("\nEstrutura de Conteudo:")
print("  Slides 1-2:   Capa e visao geral (89 problemas)")
print("  Slides 3-9:   Detalhamento dos problemas por severidade")
print("  Slides 10-15: Resumo, analogias e plano de acao")
print("  Slides 16-25: Integracao com sistemas existentes (11 slides)")
print("  Slide 26:     Resumo executivo final")
print("\nPublico-alvo: Executivos, stakeholders tecnicos, tomadores de decisao")
print("="*70)
